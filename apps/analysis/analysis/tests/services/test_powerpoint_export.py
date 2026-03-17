from io import BytesIO
from typing import Any, cast

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pydantic import ValidationError

from analysis.services.powerpoint_export import build_presentation
from analysis.web.api.schemas.datasets import PowerPointExportRequest


def test_build_presentation_for_mean_bar_uses_numeric_axis_format() -> None:
    """Mean bar exports should keep a numeric axis instead of percentages."""
    payload = {
        "file_name": "mean-bar-2026-03-17.pptx",
        "title": "Satisfaction",
        "meta_line": "Dataset: Survey 2026 | Exported: Mar 17, 2026",
        "palette": ["#3b82f6"],
        "chart": {
            "kind": "meanBar",
            "points": [
                {"label": "Mean", "value": 4.2},
                {"label": "Median", "value": 4.0},
            ],
            "min_value": 1,
            "max_value": 5,
            "color": "#3b82f6",
        },
    }

    presentation_bytes = build_presentation(payload)
    presentation = Presentation(BytesIO(presentation_bytes))

    chart = cast(Any, presentation.slides[0].shapes[2]).chart
    assert chart.value_axis.tick_labels.number_format == "0.0"


def test_build_presentation_for_distribution_preserves_point_colors() -> None:
    """Distribution exports should keep each point color in the chart."""
    payload = {
        "file_name": "bar-export-2026-03-17.pptx",
        "title": "Age group",
        "meta_line": "Dataset: Survey 2026 | Exported: Mar 17, 2026",
        "palette": ["#ff0000", "#00ff00"],
        "chart": {
            "kind": "bar",
            "points": [
                {"label": "18-29", "value": 55, "color": "#ff0000"},
                {"label": "30-44", "value": 45, "color": "#00ff00"},
            ],
        },
    }

    presentation_bytes = build_presentation(payload)
    presentation = Presentation(BytesIO(presentation_bytes))

    chart = cast(Any, presentation.slides[0].shapes[2]).chart
    assert chart.series[0].points[0].format.fill.fore_color.rgb == RGBColor.from_string(
        "FF0000"
    )
    assert chart.series[0].points[1].format.fill.fore_color.rgb == RGBColor.from_string(
        "00FF00"
    )


def test_build_presentation_for_stacked_chart_uses_contrast_label_color() -> None:
    """Stacked chart labels should use a contrasting color against segment fills."""
    payload = {
        "file_name": "stacked-export-2026-03-17.pptx",
        "title": "Awareness by segment",
        "meta_line": "Dataset: Survey 2026 | Exported: Mar 17, 2026",
        "palette": ["#1f2937", "#e5e7eb"],
        "chart": {
            "kind": "horizontalStackedBar",
            "rows": [
                {
                    "label": "Group A",
                    "segments": [
                        {"label": "Yes", "value": 60, "color": "#1f2937"},
                        {"label": "No", "value": 40, "color": "#e5e7eb"},
                    ],
                }
            ],
        },
    }

    presentation_bytes = build_presentation(payload)
    presentation = Presentation(BytesIO(presentation_bytes))

    chart = cast(Any, presentation.slides[0].shapes[2]).chart
    assert chart.series[0].points[0].data_label.font.color.rgb == RGBColor.from_string(
        "FFFFFF"
    )
    assert chart.series[1].points[0].data_label.font.color.rgb == RGBColor.from_string(
        "000000"
    )


def test_powerpoint_request_rejects_too_many_metric_cards() -> None:
    """Metrics export is limited to six cards to match the slide layout."""
    with pytest.raises(ValidationError):
        PowerPointExportRequest.model_validate(
            {
                "file_name": "metrics-2026-03-17.pptx",
                "title": "Metrics",
                "meta_line": "Dataset: Survey 2026 | Exported: Mar 17, 2026",
                "palette": ["#3b82f6"],
                "chart": {
                    "kind": "metrics",
                    "metrics": [
                        {"label": f"Metric {index}", "value": str(index)}
                        for index in range(7)
                    ],
                },
            }
        )


def test_powerpoint_request_rejects_inconsistent_stacked_rows() -> None:
    """Stacked chart rows must share the same segment order."""
    with pytest.raises(ValidationError):
        PowerPointExportRequest.model_validate(
            {
                "file_name": "stacked-2026-03-17.pptx",
                "title": "Age by gender",
                "meta_line": "Dataset: Survey 2026 | Split: Gender | Exported: Mar 17, 2026",
                "palette": ["#3b82f6", "#ef4444"],
                "chart": {
                    "kind": "horizontalStackedBar",
                    "rows": [
                        {
                            "label": "Women",
                            "segments": [
                                {"label": "Yes", "value": 60, "color": "#3b82f6"},
                                {"label": "No", "value": 40, "color": "#ef4444"},
                            ],
                        },
                        {
                            "label": "Men",
                            "segments": [
                                {"label": "No", "value": 45, "color": "#ef4444"},
                                {"label": "Yes", "value": 55, "color": "#3b82f6"},
                            ],
                        },
                    ],
                },
            }
        )


def test_powerpoint_request_rejects_inconsistent_stacked_row_colors() -> None:
    """Stacked chart rows must share the same colors for matching segments."""
    with pytest.raises(ValidationError):
        PowerPointExportRequest.model_validate(
            {
                "file_name": "stacked-colors-2026-03-17.pptx",
                "title": "Age by gender",
                "meta_line": "Dataset: Survey 2026 | Split: Gender | Exported: Mar 17, 2026",
                "palette": ["#3b82f6", "#ef4444"],
                "chart": {
                    "kind": "horizontalStackedBar",
                    "rows": [
                        {
                            "label": "Women",
                            "segments": [
                                {"label": "Yes", "value": 60, "color": "#3b82f6"},
                                {"label": "No", "value": 40, "color": "#ef4444"},
                            ],
                        },
                        {
                            "label": "Men",
                            "segments": [
                                {"label": "Yes", "value": 55, "color": "#22c55e"},
                                {"label": "No", "value": 45, "color": "#ef4444"},
                            ],
                        },
                    ],
                },
            }
        )
