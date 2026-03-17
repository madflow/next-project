# pyright: reportMissingImports=false

from io import BytesIO
from pathlib import Path
from typing import Any, Sequence, TypedDict

from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PPTX_MEDIA_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
EXPORT_FONT_SIZE_DELTA_PT = 1
TITLE_FONT_SIZE_PT = 24 + EXPORT_FONT_SIZE_DELTA_PT
META_FONT_SIZE_PT = 10 + EXPORT_FONT_SIZE_DELTA_PT
AXIS_FONT_SIZE_PT = 11 + EXPORT_FONT_SIZE_DELTA_PT
DATA_LABEL_FONT_SIZE_PT = 10 + EXPORT_FONT_SIZE_DELTA_PT
METRIC_LABEL_FONT_SIZE_PT = 10 + EXPORT_FONT_SIZE_DELTA_PT
METRIC_VALUE_FONT_SIZE_PT = 18 + EXPORT_FONT_SIZE_DELTA_PT
DEFAULT_LIGHT_TEXT_HEX = "#FFFFFF"
DEFAULT_DARK_TEXT_HEX = "#000000"
SLIDE_BACKGROUND_HEX = "#FFFFFF"


class ExportPoint(TypedDict):
    label: str
    value: float
    color: str


class ExportStackedSegment(TypedDict):
    label: str
    value: float
    color: str


class ExportStackedRow(TypedDict):
    label: str
    segments: list[ExportStackedSegment]


class ExportMetric(TypedDict):
    label: str
    value: str


def _normalize_hex_color(color: str) -> str:
    normalized = color.strip().removeprefix("#")
    if len(normalized) != 6:
        raise ValueError(f"Unsupported color value '{color}'")
    return normalized.upper()


def _rgb_components(color: str) -> tuple[int, int, int]:
    normalized = _normalize_hex_color(color)
    return (
        int(normalized[0:2], 16),
        int(normalized[2:4], 16),
        int(normalized[4:6], 16),
    )


def _hex_to_rgb_color(color: str) -> RGBColor:
    return RGBColor.from_string(_normalize_hex_color(color))


def _relative_luminance(color: str) -> float:
    def to_linear(channel: int) -> float:
        value = channel / 255
        if value <= 0.03928:
            return value / 12.92
        return ((value + 0.055) / 1.055) ** 2.4

    red, green, blue = _rgb_components(color)
    return (
        0.2126 * to_linear(red) + 0.7152 * to_linear(green) + 0.0722 * to_linear(blue)
    )


def _contrast_ratio(first_color: str, second_color: str) -> float:
    first_luminance = _relative_luminance(first_color)
    second_luminance = _relative_luminance(second_color)
    lighter = max(first_luminance, second_luminance)
    darker = min(first_luminance, second_luminance)
    return (lighter + 0.05) / (darker + 0.05)


def _best_text_color(background_color: str | None) -> RGBColor:
    if not background_color:
        return _hex_to_rgb_color(DEFAULT_LIGHT_TEXT_HEX)

    try:
        dark_contrast = _contrast_ratio(background_color, DEFAULT_DARK_TEXT_HEX)
        light_contrast = _contrast_ratio(background_color, DEFAULT_LIGHT_TEXT_HEX)
    except ValueError:
        return _hex_to_rgb_color(DEFAULT_LIGHT_TEXT_HEX)

    if dark_contrast >= light_contrast:
        return _hex_to_rgb_color(DEFAULT_DARK_TEXT_HEX)

    return _hex_to_rgb_color(DEFAULT_LIGHT_TEXT_HEX)


def _set_chart_title(slide: Any, title: str, meta_line: str) -> None:
    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.55)
    )
    title_frame = title_box.text_frame
    title_frame.clear()
    title_frame.word_wrap = True
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = title
    title_paragraph.font.size = Pt(TITLE_FONT_SIZE_PT)
    title_paragraph.font.bold = True

    meta_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.95), Inches(12.1), Inches(0.35)
    )
    meta_frame = meta_box.text_frame
    meta_frame.clear()
    meta_frame.word_wrap = True
    meta_paragraph = meta_frame.paragraphs[0]
    meta_paragraph.text = meta_line
    meta_paragraph.font.size = Pt(META_FONT_SIZE_PT)


def _apply_series_fill(series: Any, color: str) -> None:
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb_color(color)


def _apply_point_fills(series: Any, points: Sequence[ExportPoint]) -> None:
    default_color = points[0]["color"] if points else "#3b82f6"
    _apply_series_fill(series, default_color)

    for index, point in enumerate(points):
        fill = series.points[index].format.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb_color(point["color"])


def _style_data_label(data_label: Any, text_color: RGBColor) -> None:
    font = data_label.font
    font.size = Pt(DATA_LABEL_FONT_SIZE_PT)
    font.color.rgb = text_color


def _style_value_axis(
    chart: Any,
    minimum: float | None = None,
    maximum: float | None = None,
    number_format: str = '0"%"',
) -> None:
    value_axis = chart.value_axis
    value_axis.minimum_scale = minimum
    value_axis.maximum_scale = maximum
    tick_labels = value_axis.tick_labels
    tick_labels.number_format = number_format
    tick_labels.number_format_is_linked = False
    tick_labels.font.size = Pt(AXIS_FONT_SIZE_PT)


def _style_category_axis(chart: Any, reverse_order: bool = False) -> None:
    category_axis = chart.category_axis
    category_axis.reverse_order = reverse_order
    category_axis.tick_labels.font.size = Pt(AXIS_FONT_SIZE_PT)


def _configure_default_legend(chart: Any) -> None:
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False


def _configure_default_data_labels(
    plot: Any,
    *,
    position: Any,
    number_format: str = '0"%"',
    text_color: RGBColor | None = None,
) -> None:
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.position = position
    data_labels.number_format = number_format
    data_labels.number_format_is_linked = False
    data_labels.show_value = True
    data_labels.font.size = Pt(DATA_LABEL_FONT_SIZE_PT)
    if text_color is not None:
        data_labels.font.color.rgb = text_color


def _add_distribution_chart(
    slide: Any,
    chart_type: Any,
    points: Sequence[ExportPoint],
    *,
    reverse_order: bool = False,
) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = [point["label"] for point in points]
    chart_data.add_series("Series 1", [point["value"] for point in points])

    chart = slide.shapes.add_chart(
        chart_type, Inches(0.7), Inches(1.45), Inches(11.8), Inches(5.1), chart_data
    ).chart
    _style_value_axis(chart, 0.0, 100.0)
    _style_category_axis(chart, reverse_order=reverse_order)
    _apply_point_fills(chart.series[0], points)
    _configure_default_data_labels(
        chart.plots[0],
        position=XL_LABEL_POSITION.RIGHT
        if reverse_order
        else XL_LABEL_POSITION.OUTSIDE_END,
        text_color=_best_text_color(SLIDE_BACKGROUND_HEX),
    )


def _add_stacked_bar_chart(slide: Any, rows: Sequence[ExportStackedRow]) -> None:
    if not rows:
        raise ValueError("Stacked chart export requires at least one row")

    first_row = rows[0]
    chart_data = CategoryChartData()
    chart_data.categories = [row["label"] for row in rows]

    for segment_index, segment in enumerate(first_row["segments"]):
        values = [row["segments"][segment_index]["value"] for row in rows]
        chart_data.add_series(segment["label"], values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED,
        Inches(0.7),
        Inches(1.45),
        Inches(11.8),
        Inches(5.1),
        chart_data,
    ).chart
    _style_value_axis(chart, 0.0, 100.0)
    _style_category_axis(chart, reverse_order=True)
    _configure_default_legend(chart)
    _configure_default_data_labels(chart.plots[0], position=XL_LABEL_POSITION.CENTER)

    for series_index, series in enumerate(chart.series):
        _apply_series_fill(series, first_row["segments"][series_index]["color"])
        text_color = _best_text_color(first_row["segments"][series_index]["color"])
        for point in series.points:
            _style_data_label(point.data_label, text_color)


def _add_pie_chart(slide: Any, points: Sequence[ExportPoint]) -> None:
    chart_data = ChartData()
    chart_data.categories = [point["label"] for point in points]
    chart_data.add_series("Series 1", [point["value"] for point in points])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(0.9),
        Inches(1.45),
        Inches(11.2),
        Inches(5.1),
        chart_data,
    ).chart
    _configure_default_legend(chart)
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    data_labels.show_percentage = True
    data_labels.number_format = "0%"
    data_labels.number_format_is_linked = False
    data_labels.font.size = Pt(DATA_LABEL_FONT_SIZE_PT)
    data_labels.font.color.rgb = _best_text_color(SLIDE_BACKGROUND_HEX)

    for index, point in enumerate(points):
        fill = chart.series[0].points[index].format.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb_color(point["color"])


def _add_mean_bar_chart(
    slide: Any,
    points: Sequence[dict[str, float | str]],
    color: str,
    minimum: float,
    maximum: float,
) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = [str(point["label"]) for point in points]
    chart_data.add_series("Series 1", [float(point["value"]) for point in points])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.7),
        Inches(1.45),
        Inches(11.8),
        Inches(5.1),
        chart_data,
    ).chart
    _style_value_axis(chart, float(minimum), float(maximum), number_format="0.0")
    _style_category_axis(chart, reverse_order=True)
    _apply_series_fill(chart.series[0], color)

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.position = XL_LABEL_POSITION.RIGHT
    data_labels.number_format = "0.0"
    data_labels.number_format_is_linked = False
    data_labels.show_value = True
    data_labels.font.size = Pt(DATA_LABEL_FONT_SIZE_PT)
    data_labels.font.color.rgb = _best_text_color(SLIDE_BACKGROUND_HEX)


def _add_metrics_cards(
    slide: Any, metrics: Sequence[ExportMetric], palette: Sequence[str]
) -> None:
    safe_palette = list(palette) or ["#3b82f6"]
    card_width = Inches(3.8)
    card_height = Inches(1.15)
    left_positions = [Inches(0.7), Inches(4.55), Inches(8.4)]
    top_positions = [Inches(1.55), Inches(2.9)]

    for index, metric in enumerate(metrics):
        top = top_positions[index // 3]
        left = left_positions[index % 3]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_width, card_height
        )
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb_color("#ffffff")
        line = shape.line
        line.color.rgb = _hex_to_rgb_color(safe_palette[index % len(safe_palette)])
        line.width = Pt(2)

        frame = shape.text_frame
        frame.clear()
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        title_paragraph = frame.paragraphs[0]
        title_paragraph.text = metric["label"]
        title_paragraph.font.size = Pt(METRIC_LABEL_FONT_SIZE_PT)
        title_paragraph.alignment = PP_ALIGN.CENTER

        value_paragraph = frame.add_paragraph()
        value_paragraph.text = metric["value"]
        value_paragraph.font.size = Pt(METRIC_VALUE_FONT_SIZE_PT)
        value_paragraph.font.bold = True
        value_paragraph.alignment = PP_ALIGN.CENTER


def build_presentation(payload: dict[str, Any]) -> bytes:
    """Build a PowerPoint presentation from the normalized export payload."""
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _set_chart_title(slide, payload["title"], payload["meta_line"])

    chart = payload["chart"]
    kind = chart["kind"]

    if kind == "bar":
        _add_distribution_chart(slide, XL_CHART_TYPE.COLUMN_CLUSTERED, chart["points"])
    elif kind in {"horizontalBar", "multiResponse"}:
        _add_distribution_chart(
            slide, XL_CHART_TYPE.BAR_CLUSTERED, chart["points"], reverse_order=True
        )
    elif kind == "horizontalStackedBar":
        _add_stacked_bar_chart(slide, chart["rows"])
    elif kind == "pie":
        _add_pie_chart(slide, chart["points"])
    elif kind == "meanBar":
        _add_mean_bar_chart(
            slide,
            chart["points"],
            chart["color"],
            chart["min_value"],
            chart["max_value"],
        )
    elif kind == "metrics":
        _add_metrics_cards(slide, chart["metrics"], payload["palette"])
    else:
        raise ValueError(f"Unsupported chart kind '{kind}'")

    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def build_content_disposition(file_name: str) -> str:
    """Create a safe attachment header value for the generated presentation."""
    safe_name = (
        Path(file_name).name.replace("\r", "").replace("\n", "").replace('"', "")
    )
    if not safe_name:
        safe_name = "chart-export.pptx"
    return f'attachment; filename="{safe_name}"'
