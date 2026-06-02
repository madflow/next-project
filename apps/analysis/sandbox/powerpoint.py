"""Test python-pptx chart generation with various chart types and styling."""

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.util import Inches


def main() -> None:
    # create presentation with 1 slide ------
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # define chart data ---------------------
    chart_data = CategoryChartData()

    chart_data.categories = ["East", "West", "Midwest"]
    chart_data.add_series("Series 1", (19.2, 21.4, 16.7))

    # add chart to slide --------------------
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.position = XL_LABEL_POSITION.RIGHT

    category_axis = chart.category_axis
    category_axis.reverse_order = True

    prs.save("charts.pptx")


if __name__ == "__main__":
    main()
